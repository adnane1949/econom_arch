from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # ---------------------------
    # Slide 1: Title Slide
    # ---------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Plateforme Big Data : Architecture de Données"
    subtitle.text = "Collecte, Traitement et Analyse d'Articles de Presse\nSoutenance de Projet"
    
    # ---------------------------
    # Slide 2: Contexte et Objectifs
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Contexte & Objectifs du Projet"
    
    tf = body_shape.text_frame
    tf.text = "Contexte : Volumes massifs de publications médias chaque jour."
    
    p = tf.add_paragraph()
    p.text = "Objectifs (Cahier des charges) :"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Collecte Web (Scraping de sites d'actualité)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Architecture Distribuée & Big Data (Kafka, MinIO, Spark)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Data Lake avec Architecture Médaillon (Bronze, Silver, Gold)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Mise en place d'un entrepôt de données (Data Warehouse) pour l'analyse."
    p.level = 2

    # ---------------------------
    # Slide 3: Choix Technologiques
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Écosystème Technologique"
    
    tf = body_shape.text_frame
    tf.text = "Les technologies suivantes ont été déployées (via Docker) :"
    
    techs = [
        ("Python & BeautifulSoup", "Pour le Web Scraping asynchrone."),
        ("Apache Kafka", "Ingestion des événements en temps réel (Streaming)."),
        ("MinIO (S3-compatible)", "Stockage distribué du Data Lake (JSON/JSONL)."),
        ("Apache Spark", "Moteur de calcul distribué pour l'agrégation des données."),
        ("PostgreSQL", "SGBD pour l'entrepôt de données (Data Warehouse)."),
        ("Apache Airflow", "Orchestrateur central (DAG) pour automatiser le pipeline horaire.")
    ]
    for tech, desc in techs:
        p = tf.add_paragraph()
        p.text = f"{tech} : {desc}"
        p.level = 1

    # ---------------------------
    # Slide 4: Architecture Médaillon (Le Data Lake)
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "L'Architecture Médaillon"
    
    tf = body_shape.text_frame
    tf.text = "Structuration des données en trois couches de qualité croissante :"
    
    p = tf.add_paragraph()
    p.text = "1. Couche Bronze (Données Brutes) :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Historique immuable. Données stockées telles que collectées (JSON Lines)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. Couche Silver (Données Nettoyées) :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Filtrage HTML, déduplication, validation, et aplatissement (flattening)."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "3. Couche Gold (Données Analytiques) :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Agrégations via Apache Spark : Tendances, sentiments, KPIs."
    p.level = 2

    # ---------------------------
    # Slide 5: L'Orchestration et Le Pipeline
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Le Pipeline de Données (ETL / ELT)"
    
    tf = body_shape.text_frame
    tf.text = "Le workflow est orchestré par Apache Airflow (exécution horaire) :"
    
    p = tf.add_paragraph()
    p.text = "Étape 1 : Scraper vers Kafka."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 2 : Vidage Kafka vers la zone Bronze (Micro-batching)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 3 : Nettoyage ETL (Bronze -> Silver)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 4 : Job PySpark (Silver -> Gold)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 5 : Chargement des métriques Gold vers PostgreSQL."
    p.level = 1

    # ---------------------------
    # Slide 6: Qualité et Gouvernance des Données
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Qualité et Traçabilité (Gouvernance)"
    
    tf = body_shape.text_frame
    tf.text = "Afin de garantir la fiabilité des tableaux de bord finaux :"
    
    p = tf.add_paragraph()
    p.text = "Contrôles de Qualité (Complétude, Cohérence, Validité) :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Vérification des titres manquants, contenus trop courts."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Validation des intervalles de scores de sentiment [-1, 1]."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Traçabilité (Data Lineage) :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Chaque donnée est tracée avec l'identifiant du pipeline, l'URI source et la date."
    p.level = 2

    # ---------------------------
    # Slide 7: Conclusion
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Conclusion"
    
    tf = body_shape.text_frame
    tf.text = "Bilan du Projet :"
    
    p = tf.add_paragraph()
    p.text = "Architecture bout-en-bout fonctionnelle et scalable."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Intégration réussie des paradigmes Batch et Streaming."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Des données fiables prêtes pour la visualisation (Dashboards BI)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Perspectives :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Intégration d'un modèle NLP (Deep Learning) pour la détection de Fake News."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Déploiement de l'infrastructure sur un cluster Kubernetes."
    p.level = 2
    
    prs.save("presentation_projet.pptx")

if __name__ == "__main__":
    create_presentation()
    print("Presentation created successfully as 'presentation_projet.pptx'")
